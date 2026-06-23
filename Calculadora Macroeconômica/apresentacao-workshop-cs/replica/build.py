#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Monta o espelho PPTX do Workshop CS | Calculadora Macroeconomica."""
import sys, os, json

SKILL_SCRIPTS = "/sessions/zealous-dazzling-allen/mnt/.remote-plugins/plugin_01DXQJzML9C3UeVrUuhJUyC7/skills/replicar-pptx/scripts"
sys.path.insert(0, os.path.dirname(SKILL_SCRIPTS))
sys.path.insert(0, SKILL_SCRIPTS)

from coords import Page, carregar_template, salvar, layout_em_branco, ajustar_picture, cover_picture, recortar_ref
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
REF = os.path.join(BASE, "ref")
ASSETS = os.path.join(BASE, "assets-png")
LOGOS = os.path.join(BASE, "logos")
PPTX_PATH = os.path.join(BASE, "apresentacao_v2.pptx")

NAVY = RGBColor(0x00, 0x2B, 0x49)
SLATE = RGBColor(0x5E, 0x8A, 0xB4)
SIGNAL = RGBColor(0x00, 0x85, 0xCA)
GRAY1 = RGBColor(0x64, 0x64, 0x64)
GRAY2 = RGBColor(0x96, 0x96, 0x96)
MIST = RGBColor(0xD9, 0xD9, 0xD9)
STRIPE = RGBColor(0xDC, 0xE6, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTGRAY = RGBColor(0xF2, 0xF4, 0xF7)

pg = Page(w_emu=18288000, h_emu=10287000)  # tamanho real do template (20 x 11.25 in)
prs = carregar_template(PPTX_PATH)
# Forca o layout "Empty" puro (sem logo/regua baked-in no master). "Empty with
# logo" e "[Dark] Empty with logo" ja trazem um picture "Logo A&M" e uma linha
# "Linha Superior" no proprio layout, que duplicavam o wordmark/regua que o
# codigo abaixo desenha manualmente (chrome_claro / separadores).
BLANK = layout_em_branco(prs, nomes=("Empty",))

# O guia de tamanhos (espelho.md) foi calibrado para o slide padrao 16:9 de
# 13.333 x 7.5 in (12192000 EMU de largura). O template real desta apresentacao
# e 1.5x maior fisicamente (18288000 EMU), mesma proporcao. Sem compensar isso,
# todo texto sai ~1/3 menor que a referencia e a hierarquia (titulo, corpo,
# takeaway, rotulo) perde a padronizacao. Escalamos toda fonte por este fator.
FONT_SCALE = (18288000 / 12192000) * 1.19  # 1.5 base x correcao calibrada por pixel contra ref/slide-03.png


def new_slide():
    return prs.slides.add_slide(BLANK)


def no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def rect(slide, fx, fy, fw, fh, color, line=False):
    left, top, w, h = pg.box(fx, fy, fw, fh)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = color
        shp.line.width = Pt(0.5)
    no_shadow(shp)
    return shp


def line(slide, fx, fy, fw, color, thickness_pt=0.75):
    left, top, w, h = pg.box(fx, fy, fw, 0.001)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Pt(thickness_pt * FONT_SCALE))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    no_shadow(shp)
    return shp


def vline(slide, fx, fy, fh, color, thickness_pt=0.75):
    left, top, w, h = pg.box(fx, fy, 0.001, fh)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(thickness_pt * FONT_SCALE), h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    no_shadow(shp)
    return shp


def textbox(slide, fx, fy, fw, fh, runs, size=14, color=NAVY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, caps=False,
            line_spacing=1.08, italic=False, wrap=True):
    """runs: string unica OU lista de strings (uma por paragrafo, ex: bullets)."""
    left, top, w, h = pg.box(fx, fy, fw, fh)
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    from pptx.enum.text import MSO_AUTO_SIZE
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paragraphs = runs if isinstance(runs, list) else [runs]
    for i, ptext in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(ptext, tuple):
            txt, kwargs = ptext
        else:
            txt, kwargs = ptext, {}
        run = p.add_run()
        run.text = txt.upper() if caps else txt
        run.font.size = Pt(kwargs.get("size", size) * FONT_SCALE)
        run.font.bold = kwargs.get("bold", bold)
        run.font.italic = kwargs.get("italic", italic)
        run.font.color.rgb = kwargs.get("color", color)
    return box


def picture(slide, fx, fy, fw, fh, path):
    left, top, w, h = pg.box(fx, fy, fw, fh)
    return ajustar_picture(slide, path, left, top, w, h)


def fullbleed(slide, path):
    left, top, w, h = pg.box(0, 0, 1, 1)
    return cover_picture(slide, path, w, h)


def chrome_claro(slide, eyebrow_txt, titulo_txt, pagina_txt, eyebrow_w=0.55):
    # posicoes calibradas em pixel contra ref/slide-03.png (eyebrow ~0.027,
    # regua ~0.052, titulo ~0.090 de altura) em vez das fracoes genericas do
    # espelho.md, que estavam mais espacadas que a referencia real.
    if eyebrow_txt:
        textbox(slide, 0.04, 0.016, eyebrow_w, 0.025, eyebrow_txt, size=10.5,
                color=SLATE, bold=False, caps=True)
    picture(slide, 0.855, 0.012, 0.105, 0.03, os.path.join(LOGOS, "am-wordmark-navy.png"))
    line(slide, 0.04, 0.052, 0.92, NAVY)
    textbox(slide, 0.04, 0.071, 0.70, 0.075, titulo_txt, size=22, color=NAVY, bold=False)
    textbox(slide, 0.935, 0.955, 0.04, 0.03, pagina_txt, size=9.5, color=GRAY2,
            align=PP_ALIGN.RIGHT)


def takeaway(slide, fx, fy, fw, fh, txt, color=NAVY, size=13.5):
    textbox(slide, fx, fy, fw, fh, txt, size=size, color=color, bold=True, line_spacing=1.12)


def fonte(slide, fx, fy, fw, txt):
    textbox(slide, fx, fy, fw, 0.03, txt, size=8, color=GRAY2)


# ============================================================ SLIDE 01 — CAPA
s = new_slide()
fullbleed(s, os.path.join(ASSETS, "covers_tecnologia.png"))
picture(s, 0.025, 0.035, 0.16, 0.08, os.path.join(LOGOS, "am-infra-white.png"))
line(s, 0.025, 0.585, 0.45, WHITE, 1.0)
textbox(s, 0.025, 0.605, 0.60, 0.09, "Workshop CS", size=30, color=WHITE, bold=False)
textbox(s, 0.025, 0.705, 0.65, 0.06,
        "Calculadora Macroeconômica: de repositório a apresentação",
        size=14, color=WHITE)
textbox(s, 0.80, 0.915, 0.17, 0.04, "22.06.2026", size=11, color=WHITE, align=PP_ALIGN.RIGHT)

# ============================================================ SLIDE 02 — SEPARADOR 01 (escuro)
s = new_slide()
rect(s, 0, 0, 1, 1, NAVY)
line(s, 0.0, 0.0, 1.0, RGBColor(0x3D, 0x5A, 0x73), 1.5)
picture(s, 0.855, 0.022, 0.115, 0.035, os.path.join(LOGOS, "am-wordmark-white.png"))
textbox(s, 0.02, 0.30, 0.20, 0.18, "01", size=64, color=SLATE)
textbox(s, 0.02, 0.465, 0.60, 0.11, "Da dor ao deck em minutos", size=27, color=WHITE, bold=False)

# ============================================================ SLIDE 03 — Objetivos deste workshop
s = new_slide()
chrome_claro(s, "DA DOR AO DECK EM MINUTOS", "Objetivos deste workshop", "1")

itens = [
    ("01", "decoracao_icones_search-slate.png",
     "Como a pipeline lê um projeto real (código, dados, documentação) e gera um roteiro estruturado"),
    ("02", "decoracao_icones_shield-check-slate.png",
     "Como cada etapa (roteiro, layout, gráficos, decoração, código, QA) é revisada por um agente especialista antes de avançar"),
    ("03", "decoracao_icones_file-text-slate.png",
     "Este próprio deck como prova: foi gerado a partir da pasta do projeto Calculadora Macroeconômica"),
]
y0 = 0.275
row_h = 0.105
for i, (num, icon, txt) in enumerate(itens):
    y = y0 + i * row_h
    textbox(s, 0.02, y, 0.035, 0.04, num, size=13, color=SIGNAL, bold=True)
    picture(s, 0.058, y + 0.003, 0.028, 0.032, os.path.join(ASSETS, icon))
    textbox(s, 0.10, y - 0.005, 0.39, 0.09, txt, size=11.5, color=NAVY, line_spacing=1.12)
    if i < 2:
        line(s, 0.02, y + row_h - 0.018, 0.47, MIST)

rect(s, 0.50, 0.255, 0.48, 0.095, LIGHTGRAY)
textbox(s, 0.52, 0.27, 0.40, 0.05, "7 etapas", size=24, color=NAVY, bold=True)
textbox(s, 0.52, 0.325, 0.44, 0.03, "do roteiro ao PDF final", size=10.5, color=GRAY1)
line(s, 0.50, 0.365, 0.45, MIST)
textbox(s, 0.52, 0.40, 0.45, 0.05,
        "Propostas e relatórios de Capital Strategy nascem de planilhas, modelos e memos densos",
        size=11.5, color=NAVY)
textbox(s, 0.52, 0.465, 0.45, 0.04,
        "Hoje, tempo de formatação compete com tempo de análise", size=11.5, color=NAVY)

takeaway(s, 0.02, 0.865, 0.93, 0.09,
         "Ao final deste workshop, o time de Capital Strategy entende como transformar um repositório "
         "de código e dados em uma apresentação pronta, sem depender de designer ou de horas de "
         "formatação manual.")

# ============================================================ SLIDE 04 — A dor
s = new_slide()
chrome_claro(s, "DA DOR AO DECK EM MINUTOS", "A dor: dados prontos, apresentação represada", "2")

textbox(s, 0.02, 0.235, 0.30, 0.04, "Sintomas comuns", size=13, color=NAVY, bold=True)
sintomas = [
    "Horas redistribuindo a mesma tabela em formatos diferentes a cada apresentação",
    "Gráficos refeitos manualmente sempre que um dado é atualizado",
    "Identidade visual inconsistente entre decks do mesmo projeto",
    "Revisão de conteúdo e revisão de design acontecendo nas mesmas idas e voltas",
]
textbox(s, 0.02, 0.285, 0.52, 0.30, ["•  " + t for t in sintomas], size=12, color=NAVY, line_spacing=1.22)

textbox(s, 0.61, 0.235, 0.34, 0.06, "8–12", size=30, color=SIGNAL, bold=True)
textbox(s, 0.61, 0.30, 0.34, 0.04, "slides por proposta A&M, hoje gastos em formatação",
        size=10.5, color=GRAY1)
rect(s, 0.61, 0.395, 0.335, 0.225, WHITE, line=True)
textbox(s, 0.635, 0.415, 0.29, 0.04, "Custo real", size=12.5, color=NAVY, bold=True)
textbox(s, 0.635, 0.46, 0.29, 0.15,
        ["Cada deck consome ciclos de analista sênior em tarefa de formatação, não de julgamento.",
         "É o intervalo que, hoje, amarra contexto, escopo e investimento em cada proposta."],
        size=11, color=NAVY, line_spacing=1.18)

takeaway(s, 0.02, 0.84, 0.82, 0.08,
         "O gargalo de uma apresentação raramente é a análise, é transformar a análise em slides "
         "consistentes e bem desenhados na velocidade que o negócio pede.")
fonte(s, 0.02, 0.945, 0.5, "Fonte: referência interna A&M, estrutura-deck.md")

# ============================================================ SLIDE 05 — Como funciona
s = new_slide()
chrome_claro(s, "DA DOR AO DECK EM MINUTOS", "Como funciona o apresentacao-rapida-am", "3")
picture(s, 0.02, 0.24, 0.96, 0.16, os.path.join(ASSETS, "decoracao__valida_d1.png"))

textbox(s, 0.02, 0.455, 0.46, 0.035, "O que isso garante", size=12.5, color=NAVY, bold=True)
textbox(s, 0.02, 0.495, 0.46, 0.15,
        "Cada uma das 7 etapas, do roteiro ao código LaTeX, passa por validação própria antes de "
        "avançar para a próxima, o que elimina retrabalho de formatação e inconsistência visual "
        "entre seções do mesmo deck.", size=11.5, color=NAVY, line_spacing=1.2)

textbox(s, 0.50, 0.455, 0.46, 0.035, "O que isso elimina", size=12.5, color=NAVY, bold=True)
textbox(s, 0.52, 0.495, 0.44, 0.10,
        ["•  Retrabalho manual de formatação e diagramação",
         "•  Inconsistência visual entre seções do mesmo deck"],
        size=11.5, color=NAVY, line_spacing=1.2)

takeaway(s, 0.02, 0.84, 0.82, 0.09,
         "Dividir a criação do deck em sete etapas com validação cruzada entre agentes é o que "
         "permite trocar reunião de design por revisão automática, sem perder rigor.")

# ============================================================ SLIDE 06 — SEPARADOR 02 (claro)
s = new_slide()
rect(s, 0, 0, 1, 1, WHITE)
line(s, 0.0, 0.0, 1.0, SIGNAL, 1.5)
picture(s, 0.855, 0.022, 0.115, 0.035, os.path.join(LOGOS, "am-wordmark-navy.png"))
textbox(s, 0.02, 0.30, 0.20, 0.18, "02", size=64, color=SLATE)
textbox(s, 0.02, 0.465, 0.66, 0.13,
        "Estudo de caso ao vivo: Calculadora Macroeconômica", size=27, color=NAVY, bold=False,
        line_spacing=1.08)

# ============================================================ SLIDE 07 — Por dentro do projeto
s = new_slide()
chrome_claro(s, "ESTUDO DE CASO AO VIVO: CALCULADORA MACROECONÔMICA",
             "Por dentro do projeto: de coletores a interface", "4", eyebrow_w=0.65)

cards = [
    ("decoracao_icones_database-navy.png", "Coletores",
     "BCB SGS (séries realizadas), Boletim Focus (projeções), ANBIMA ETTJ (inflação implícita), "
     "B3 DI1 (curva de juros) e IBGE Sidra."),
    ("decoracao_icones_cog-navy.png", "Motor",
     "Módulos deflator, index_builder, projector e vector calculam fator acumulado e vetores mensais."),
    ("decoracao_icones_monitor-navy.png", "Interface",
     "App Streamlit com 6 abas: Calculadora, Vetor de Correção, ETTJ, Curva de Juros, "
     "Log de Atualizações, Documentação."),
]
xs = [0.02, 0.34, 0.66]
for x, (icon, label, txt) in zip(xs, cards):
    rect(s, x, 0.255, 0.30, 0.225, WHITE, line=True)
    picture(s, x + 0.02, 0.285, 0.028, 0.032, os.path.join(ASSETS, icon))
    textbox(s, x + 0.06, 0.28, 0.22, 0.04, label, size=13, color=NAVY, bold=True)
    textbox(s, x + 0.02, 0.335, 0.26, 0.13, txt, size=10.5, color=NAVY, line_spacing=1.18)

stats = [("7", "índices suportados"), ("4", "fontes de dados"), ("1–7 dias", "ciclo de cache")]
for x, (num, label) in zip(xs, stats):
    textbox(s, x, 0.575, 0.30, 0.06, num, size=24, color=SIGNAL, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x, 0.645, 0.30, 0.03, label, size=10, color=GRAY1, align=PP_ALIGN.CENTER)
vline(s, 0.335, 0.58, 0.10, MIST)
vline(s, 0.655, 0.58, 0.10, MIST)

textbox(s, 0.13, 0.735, 0.74, 0.06,
        "Um único arquivo, ENGINE_CONTEXT.md, já documenta contratos e limites do motor: código "
        "dividido por responsabilidade vira roteiro sem reunião de levantamento.",
        size=10.5, color=GRAY1, align=PP_ALIGN.CENTER, line_spacing=1.15)

takeaway(s, 0.02, 0.84, 0.78, 0.08,
         "A Calculadora Macroeconômica organiza a correção monetária em três camadas claras, "
         "captura de dados, motor de cálculo e interface, o que tornou a leitura automática do "
         "projeto direta.")
fonte(s, 0.02, 0.945, 0.5, "Fonte: ENGINE_CONTEXT.md, repositório do projeto (2026)")

# ============================================================ SLIDE 08 — Metodologia
s = new_slide()
chrome_claro(s, "ESTUDO DE CASO AO VIVO: CALCULADORA MACROECONÔMICA",
             "A metodologia que sustenta cada número", "5", eyebrow_w=0.65)

rows_data = [
    ("Índice", "Fonte SGS", "Início", "Projeção"),
    ("IPCA", "433", "07/1994", "Focus + ETTJ"),
    ("IGP-M", "189", "01/1940", "Focus"),
    ("IPCA-15", "7478", "02/1999", "Focus"),
    ("INPC", "188", "01/1979", "Focus"),
    ("INCC", "192", "01/1985", "–"),
    ("SELIC", "4390", "06/1986", "–"),
    ("CDI", "4391", "06/1986", "–"),
]
left, top, w, h = pg.box(0.02, 0.255, 0.38, 0.225)
gtable = s.shapes.add_table(len(rows_data), 4, left, top, w, h).table
try:
    gtable.first_row = False
    gtable.horz_banding = False
except Exception:
    pass
colw = [int(w * f) for f in (0.30, 0.24, 0.22, 0.24)]
for i, cw in enumerate(colw):
    gtable.columns[i].width = cw
rowh = int(h / len(rows_data))
for r, rowvals in enumerate(rows_data):
    gtable.rows[r].height = rowh
    for c, val in enumerate(rowvals):
        cell = gtable.cell(r, c)
        cell.margin_left = Emu(int(36000 * FONT_SCALE))
        cell.margin_right = Emu(int(36000 * FONT_SCALE))
        cell.margin_top = Emu(int(9000 * FONT_SCALE))
        cell.margin_bottom = Emu(int(9000 * FONT_SCALE))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        elif r % 2 == 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = STRIPE
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        tf = cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = val
        run.font.size = Pt(11 * FONT_SCALE)
        run.font.bold = (r == 0)
        run.font.color.rgb = WHITE if r == 0 else NAVY

textbox(s, 0.62, 0.255, 0.34, 0.04, "Como o fator é calculado", size=12.5, color=NAVY, bold=True)
textbox(s, 0.62, 0.30, 0.34, 0.05, "factor = ∏(1+rₜ/100)", size=15, color=SIGNAL, bold=True)
textbox(s, 0.62, 0.365, 0.34, 0.17,
        ["•  Realizado: produto encadeado das taxas mensais entre as duas datas",
         "•  Focus: taxa anual mediana distribuída em taxa mensal composta uniforme",
         "•  ETTJ: taxas forward da curva de inflação implícita ANBIMA (NTN-B), só para IPCA"],
        size=10.5, color=NAVY, line_spacing=1.18)
line(s, 0.62, 0.555, 0.30, MIST)
textbox(s, 0.62, 0.575, 0.34, 0.09,
        "IPCA, IGP-M, IPCA-15 e INPC têm projeção Focus, e o IPCA soma também a opção ETTJ. INCC, "
        "SELIC e CDI não projetam: a data final fica limitada ao último dado realizado.",
        size=10, color=GRAY1, line_spacing=1.16)

takeaway(s, 0.02, 0.84, 0.78, 0.08,
         "Cada valor corrigido combina dados realizados em fator acumulado com projeções de "
         "mercado, sempre identificando qual parte do cálculo é fato e qual é expectativa.")
fonte(s, 0.02, 0.945, 0.6, "Fonte: ENGINE_CONTEXT.md, API SGS e Boletim Focus do Banco Central do Brasil (2026)")

# ============================================================ SLIDE 09 — Recapitulando
s = new_slide()
chrome_claro(s, "ESTUDO DE CASO AO VIVO: CALCULADORA MACROECONÔMICA",
             "Recapitulando e como usar a ferramenta", "6", eyebrow_w=0.65)

recap = [
    ("01", "A pipeline lê código, documentação e dados e devolve um roteiro estruturado e revisado"),
    ("02", "Cada etapa tem um agente dedicado à qualidade"),
    ("03", "Estudo de caso real do início ao fim"),
]
y0 = 0.24
for i, (num, txt) in enumerate(recap):
    y = y0 + i * 0.058
    textbox(s, 0.02, y, 0.04, 0.04, num, size=13, color=SIGNAL, bold=True)
    textbox(s, 0.07, y, 0.58, 0.05, txt, size=12, color=NAVY)

picture(s, 0.02, 0.452, 0.03, 0.038, os.path.join(ASSETS, "decoracao_icones_rocket-navy.png"))
textbox(s, 0.06, 0.45, 0.40, 0.04, "Próximos passos", size=12.5, color=NAVY, bold=True)
textbox(s, 0.04, 0.50, 0.46, 0.06,
        "•  Leve para o apresentacao-rapida-am o próximo memo, modelo ou repositório que precisa virar deck",
        size=11, color=NAVY, line_spacing=1.18)
textbox(s, 0.04, 0.565, 0.46, 0.06,
        "•  Acesse a ferramenta em cortex.enpower.com.br/plugins/apresentacao-rapida-am",
        size=11, color=NAVY, line_spacing=1.18)

qr_dst = os.path.join(ASSETS, "qr-recorte.png")
recortar_ref(os.path.join(REF, "slide-09.png"), (0.695, 0.475, 0.135, 0.225), qr_dst)
picture(s, 0.695, 0.475, 0.135, 0.225, qr_dst)
textbox(s, 0.655, 0.715, 0.225, 0.03, "Criado com apresentacao-rapida-am", size=8.5, color=GRAY1,
        align=PP_ALIGN.CENTER)

takeaway(s, 0.02, 0.84, 0.78, 0.08,
         "Este deck saiu da pasta do projeto para o PDF final sem reunião de design, e o mesmo "
         "caminho está disponível para qualquer dor de dados do time de Capital Strategy.")

# ============================================================ SLIDE 10 — CONTRACAPA
s = new_slide()
fullbleed(s, os.path.join(ASSETS, "covers_backcover.png"))
picture(s, 0.40, 0.78, 0.20, 0.10, os.path.join(LOGOS, "am-infra-white.png"))

salvar(prs, PPTX_PATH)
print("OK: %d slides salvos em %s" % (len(prs.slides.__iter__.__self__._sldIdLst), PPTX_PATH))
# --- fim do script (linhas de buffer contra truncamento do mount) ---
#
#
